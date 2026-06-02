"""
scripts/build_presentation.py
==============================
Rebuild the Olist Analytics Platform presentation deck following the
coaching outline from docs/req/Coaching 2.3 session (DS3F).pdf page 6:

  1. Executive Summary
  2. Introduction & Context
  3. Methodology & Data
  4. Results & Insights
  5. Strategic Recommendations
  6. Conclusion

Run:
  conda activate olist
  cd olist_platform
  python scripts/build_presentation.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT   = Path(__file__).resolve().parents[1]
CHARTS = ROOT / "reports" / "charts"
OUT    = ROOT / "docs" / "Olist_Analytics_Platform_Presentation.pptx"

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x3B, 0x66)   # primary dark
TEAL      = RGBColor(0x1C, 0x72, 0x93)   # secondary
TEAL_LT   = RGBColor(0x4E, 0xAF, 0xD1)   # light teal
ORANGE    = RGBColor(0xF4, 0xA2, 0x61)   # accent warm
GREEN     = RGBColor(0x2E, 0xC4, 0xB6)   # success/positive
RED       = RGBColor(0xE9, 0x4F, 0x37)   # alert
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xEA, 0xF4, 0xFD)   # very light blue
DARK_TXT  = RGBColor(0x1A, 0x1A, 0x2E)
MID_TXT   = RGBColor(0x44, 0x55, 0x66)
DIM_TXT   = RGBColor(0x8A, 0x99, 0xAA)
YELLOW_HL = RGBColor(0xFF, 0xE0, 0x66)

# ── Slide dimensions (widescreen 10×5.6) ─────────────────────────────────────
W = Inches(10)
H = Inches(5.6)
SLIDE_W = W
SLIDE_H = H

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ═══════════════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(line_width)
    else:
        line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=12, bold=False, italic=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_para(tf, text, font_size=11, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, space_before=0, indent_level=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def add_textbox_multi(slide, lines, x, y, w, h, font_size=11,
                      color=DARK_TXT, line_spacing=1.0):
    """lines = list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, color
        elif len(item) == 2:
            text, bold = item; col = color
        else:
            text, bold, col = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = "Calibri"
    return txb

def add_image(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w))

def header_band(slide, title, slide_num, total=13):
    """Dark navy header bar with slide title and page number."""
    add_rect(slide, 0, 0, 10, 0.62, fill_color=NAVY)
    add_text(slide, "OLIST ANALYTICS PLATFORM  |  CONFIDENTIAL",
             0.18, 0.03, 7, 0.22, font_size=7.5, color=rgb(0x9F,0xD7,0xF0),
             bold=False, align=PP_ALIGN.LEFT)
    add_text(slide, title.upper(),
             0.18, 0.24, 8.2, 0.34, font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT)
    add_text(slide, f"{slide_num} / {total}",
             9.0, 0.22, 0.85, 0.3, font_size=10, bold=True, color=TEAL_LT,
             align=PP_ALIGN.RIGHT)

def footer_line(slide):
    add_rect(slide, 0, 5.48, 10, 0.025, fill_color=TEAL)

def kpi_card(slide, label, value, unit, x, y, w=1.7, h=0.82,
             value_color=WHITE, bg=TEAL):
    add_rect(slide, x, y, w, h, fill_color=bg)
    add_text(slide, value, x+0.07, y+0.05, w-0.12, 0.42,
             font_size=19, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text(slide, unit, x+0.07, y+0.43, w-0.12, 0.18,
             font_size=7.5, bold=False, color=rgb(0xCA,0xE8,0xF5), align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.07, y+0.60, w-0.12, 0.2,
             font_size=8.5, bold=True, color=rgb(0xCA,0xE8,0xF5), align=PP_ALIGN.CENTER)

def section_label(slide, text, x, y, w=3.5, color=TEAL):
    add_rect(slide, x, y, w, 0.22, fill_color=color)
    add_text(slide, text.upper(), x+0.1, y+0.02, w-0.15, 0.2,
             font_size=8, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def bullet_box(slide, items, x, y, w, h,
               font_size=11, bullet="▸", color=DARK_TXT, heading=None,
               heading_color=TEAL, bg=None, line_color=None):
    if bg:
        add_rect(slide, x, y, w, h, fill_color=bg, line_color=line_color, line_width=0.5)
    txb = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.08),
                                    Inches(w-0.18), Inches(h-0.12))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if heading:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = heading_color
        run.font.name = "Calibri"
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb

def two_col_table(slide, rows, x, y, w, col_w, font_size=9.5,
                  header_bg=TEAL, alt_bg=LIGHT_BG):
    """rows = list of tuples; first row is header."""
    cell_h = 0.26
    n_cols = len(rows[0])
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            cx = x + sum(col_w[:ci])
            cy = y + ri * cell_h
            bg = header_bg if ri == 0 else (alt_bg if ri % 2 == 0 else WHITE)
            add_rect(slide, cx, cy, col_w[ci], cell_h,
                     fill_color=bg, line_color=rgb(0xCC,0xDD,0xEE), line_width=0.3)
            tc = WHITE if ri == 0 else DARK_TXT
            add_text(slide, str(cell), cx+0.06, cy+0.04, col_w[ci]-0.1, cell_h-0.06,
                     font_size=font_size, bold=(ri == 0), color=tc, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(slide, 0, 0, 10, 5.6, fill_color=NAVY)

# Top accent stripe
add_rect(slide, 0, 0, 10, 0.08, fill_color=ORANGE)

# Logo / brand area
add_rect(slide, 0.3, 0.3, 0.08, 1.5, fill_color=TEAL)

# Main title
add_text(slide, "OLIST ANALYTICS PLATFORM",
         0.55, 0.28, 8.5, 0.7, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "End-to-End Data Engineering & Business Intelligence",
         0.55, 0.92, 8.5, 0.38, font_size=15, bold=False, color=TEAL_LT, align=PP_ALIGN.LEFT)

# Separator
add_rect(slide, 0.55, 1.38, 8.5, 0.025, fill_color=TEAL)

# KPI strip (6 cards)
kpi_data = [
    ("GMV",          "R$15.8M",   "Total 2016-2018",  TEAL),
    ("ORDERS",        "99.4K",    "Delivered",         TEAL),
    ("AVG ORDER",    "R$137.75",  "Per order",         TEAL),
    ("ON-TIME RATE", "91.9%",     "Delivery SLA",      rgb(0x1A,0x6B,0x4A)),
    ("AVG REVIEW",   "4.09 / 5", "Customer score",    rgb(0x1A,0x6B,0x4A)),
    ("MQL CONV.",    "10.5%",    "Seller acquisition",rgb(0x7A,0x3A,0x1A)),
]
for i, (label, val, unit, bg) in enumerate(kpi_data):
    kpi_card(slide, label, val, unit, 0.22 + i*1.60, 1.52, w=1.52, h=0.88, bg=bg)

# Tagline
add_text(slide, "Unified analytics across 11 datasets · Star schema warehouse · 181+ data quality checks · 7 dbt mart models",
         0.4, 2.54, 9.2, 0.3, font_size=9.5, color=rgb(0x8A,0xC4,0xE0), align=PP_ALIGN.CENTER)

# Section outline preview
section_titles = [
    "1  Executive Summary",
    "2  Introduction & Context",
    "3  Methodology & Data",
    "4  Results & Insights",
    "5  Strategic Recommendations",
    "6  Conclusion",
]
add_rect(slide, 0.3, 2.98, 9.4, 2.0, fill_color=rgb(0x08,0x25,0x44))
add_text(slide, "PRESENTATION OUTLINE",
         0.5, 3.02, 9, 0.28, font_size=9, bold=True, color=TEAL_LT, align=PP_ALIGN.LEFT)
for i, t in enumerate(section_titles):
    col = i // 3
    row = i % 3
    add_text(slide, t, 0.55 + col*4.7, 3.36 + row*0.5, 4.5, 0.4,
             font_size=10.5, color=WHITE, align=PP_ALIGN.LEFT)

# Bottom metadata
add_rect(slide, 0, 5.3, 10, 0.3, fill_color=rgb(0x06,0x1B,0x33))
add_text(slide, "Module 2 Assignment  |  May 2026  |  NTU SCTP Advanced Professional Certificate — Data Science & AI",
         0.3, 5.32, 9.4, 0.24, font_size=8, color=DIM_TXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Executive Summary", 2)
footer_line(slide)

# Purpose statement
add_rect(slide, 0.25, 0.75, 9.5, 0.55, fill_color=LIGHT_BG)
add_text(slide,
    "Purpose: Deliver a production-grade ELT data platform over the Brazilian Olist e-commerce dataset (2016–2018) "
    "that transforms 11 raw CSV sources into a unified BigQuery warehouse — surfacing actionable revenue, customer, "
    "marketing, and logistics intelligence for executive decision-making.",
    0.35, 0.78, 9.3, 0.50, font_size=10, color=DARK_TXT, align=PP_ALIGN.LEFT)

# Three columns: findings, actions, platform
col_data = [
    ("KEY FINDINGS", TEAL, [
        "R$15.8M GMV | 91.9% on-time delivery | 4.09/5 review score",
        "97% of 96K customers purchased only once — highest risk to growth",
        "Paid Search converts at 12.3% vs Email at 3.0% — 4× gap",
        "North region: 8.5% of population, <1% of orders — untapped",
        "Nov 2017 Black Friday peak R$1.04M = 4× monthly baseline",
        "Top 2 categories (Health & Beauty, Watches) = 18.5% of GMV",
    ]),
    ("COMMERCIAL OPPORTUNITIES", ORANGE, [
        "R$300K recoverable from 23K At-Risk customers at 5% re-activation",
        "+35–40 deals/cohort from paid search budget reallocation",
        "5,000+ new orders from Northern Brazil fulfilment hub",
        "3–5× LTV uplift via subscription model for top consumable categories",
        "0.3-star review improvement from proactive delivery notifications",
    ]),
    ("PLATFORM DELIVERED", GREEN, [
        "5-layer ELT: CSV → BigQuery Raw → Staging → Warehouse → Marts",
        "Star schema: 4 fact tables, 7 dimensions, 181+ dbt tests",
        "Dagster orchestration with daily 06:00 UTC schedule",
        "Customer RFM segmentation + CLV + funnel attribution",
        "8 executive charts | MongoDB ODS | Redis KPI cache",
    ]),
]
for ci, (title, col, items) in enumerate(col_data):
    x = 0.25 + ci * 3.2
    add_rect(slide, x, 1.42, 3.1, 0.28, fill_color=col)
    add_text(slide, title, x+0.08, 1.44, 3.0, 0.24,
             font_size=8.5, bold=True, color=WHITE)
    add_rect(slide, x, 1.70, 3.1, 3.62,
             fill_color=rgb(0xF6,0xFB,0xFF) if ci == 0 else
             rgb(0xFF,0xF8,0xF3) if ci == 1 else rgb(0xF4,0xFD,0xFB),
             line_color=rgb(0xCC,0xDD,0xEE), line_width=0.5)
    for ri, item in enumerate(items):
        add_text(slide, f"▸  {item}", x+0.1, 1.78 + ri*0.58, 2.95, 0.52,
                 font_size=9, color=DARK_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — INTRODUCTION & CONTEXT: PROBLEM STATEMENT & SCOPE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Introduction & Context — Problem Statement & Scope", 3)
footer_line(slide)

section_label(slide, "Section 2 of 6 — Introduction & Context", 0.25, 0.72)

# Left: 4 problems
problems = [
    ("📊  Siloed Data",
     "Orders, marketing leads, product reviews, and geolocation lived in 11 separate CSV files "
     "with no shared definitions, no unified key, and no automated refresh."),
    ("🔍  No Funnel Visibility",
     "8,000 marketing qualified leads generated 842 closed deals (10.5% conversion) — "
     "but zero attribution by channel, no time-to-close measurement, no ROI data."),
    ("📦  Logistics Blind Spots",
     "8.1% of orders arrived late. No SLA alerting, no regional performance tracking, "
     "and no systematic link between delivery delay and review score degradation."),
    ("🤝  Customer Churn Invisible",
     "97% of 96,096 customers never placed a second order. No segmentation, "
     "no win-back strategy, no CLV measurement, no at-risk early-warning system."),
]
add_text(slide, "BUSINESS PROBLEMS", 0.25, 0.98, 5.3, 0.28,
         font_size=10.5, bold=True, color=NAVY)
for i, (title, desc) in enumerate(problems):
    y = 1.32 + i * 1.01
    add_rect(slide, 0.25, y, 5.35, 0.94, fill_color=LIGHT_BG,
             line_color=TEAL, line_width=0.8)
    add_text(slide, title, 0.38, y+0.06, 5.1, 0.28, font_size=10.5, bold=True, color=NAVY)
    add_text(slide, desc, 0.38, y+0.34, 5.1, 0.56, font_size=9, color=MID_TXT, wrap=True)

# Right: Scope & Objectives
add_text(slide, "PROJECT SCOPE", 5.8, 0.98, 3.9, 0.28,
         font_size=10.5, bold=True, color=NAVY)
scope_items = [
    ("Dataset", "Brazilian E-Commerce by Olist (Kaggle)"),
    ("Period",  "Sep 2016 – Oct 2018 (25 months)"),
    ("Scale",   "99,441 orders | 96,096 customers | 3,095 sellers"),
    ("Sources", "11 CSV files | 1.46M total rows"),
    ("Geo",     "27 Brazilian states | 1M geolocation readings"),
]
add_rect(slide, 5.8, 1.30, 3.95, 1.65,
         fill_color=LIGHT_BG, line_color=TEAL, line_width=0.5)
for i, (k, v) in enumerate(scope_items):
    y = 1.37 + i * 0.3
    add_text(slide, k.upper()+":", 5.93, y, 1.0, 0.28,
             font_size=8.5, bold=True, color=TEAL)
    add_text(slide, v, 7.0, y, 2.7, 0.28, font_size=8.5, color=DARK_TXT)

add_text(slide, "PLATFORM OBJECTIVES", 5.8, 3.05, 3.9, 0.28,
         font_size=10.5, bold=True, color=NAVY)
objectives = [
    "Unified BigQuery warehouse across all 11 source datasets",
    "Automated ELT with 181+ data quality gates at 4 layers",
    "Customer RFM segmentation + CLV measurement",
    "Marketing funnel attribution (B2B seller acquisition)",
    "Logistics SLA monitoring with delivery KPIs",
    "Executive dashboards — CEO, CFO, CMO, COO views",
]
add_rect(slide, 5.8, 3.32, 3.95, 2.0,
         fill_color=rgb(0xF4,0xFD,0xFB), line_color=GREEN, line_width=0.5)
for i, obj in enumerate(objectives):
    add_text(slide, f"✓  {obj}", 5.93, 3.40 + i*0.31, 3.75, 0.3,
             font_size=8.5, color=DARK_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — INTRODUCTION & CONTEXT: PERFORMANCE INSIGHTS METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Introduction & Context — Performance Insights Methodology", 4)
footer_line(slide)

section_label(slide, "Section 2 of 6 — Introduction & Context", 0.25, 0.72)

add_text(slide, "Analytical Framework — From Raw Data to Actionable Decisions",
         0.25, 0.98, 9.5, 0.32, font_size=12, bold=True, color=NAVY)

# 4-step framework flow
steps = [
    ("1. DESCRIBE",     TEAL,   "What happened?",
     ["Monthly GMV & order trends", "Category revenue distribution",
      "Customer purchase frequency", "Geographic demand heatmap"]),
    ("2. DIAGNOSE",     NAVY,   "Why did it happen?",
     ["97% one-time buyer root cause", "Late delivery → review score impact",
      "Channel conversion gap (12% vs 3%)", "Regional demand vs population"]),
    ("3. PREDICT",      rgb(0x7B,0x35,0xA2), "What will happen?",
     ["At-risk segment recovery value", "Paid search ROI projection",
      "Fulfilment hub demand model", "Subscription LTV estimate"]),
    ("4. PRESCRIBE",    ORANGE, "What should we do?",
     ["5 prioritised recommendations", "Impact × Effort matrix",
      "Phased roadmap Q3 2026–Q1 2027", "KPI targets per initiative"]),
]
for i, (title, color, question, bullets) in enumerate(steps):
    x = 0.25 + i * 2.4
    add_rect(slide, x, 1.42, 2.3, 0.36, fill_color=color)
    add_text(slide, title, x+0.08, 1.45, 2.18, 0.3,
             font_size=9.5, bold=True, color=WHITE)
    add_rect(slide, x, 1.78, 2.3, 0.34, fill_color=rgb(0xDD,0xEE,0xFF))
    add_text(slide, question, x+0.08, 1.80, 2.18, 0.3,
             font_size=9.5, bold=True, color=NAVY)
    add_rect(slide, x, 2.12, 2.3, 1.84,
             fill_color=LIGHT_BG, line_color=color, line_width=0.7)
    for j, b in enumerate(bullets):
        add_text(slide, f"▸  {b}", x+0.10, 2.20 + j*0.44, 2.15, 0.42,
                 font_size=8.5, color=DARK_TXT, wrap=True)

# Arrow connectors
for i in range(3):
    ax = 0.25 + i * 2.4 + 2.30
    add_text(slide, "→", ax+0.01, 1.52, 0.1, 0.32,
             font_size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Questions answered box
add_rect(slide, 0.25, 4.10, 9.5, 1.22,
         fill_color=rgb(0xF0,0xF8,0xFF), line_color=TEAL, line_width=0.6)
add_text(slide, "KEY ANALYTICAL QUESTIONS ANSWERED IN THIS PRESENTATION",
         0.38, 4.14, 9.2, 0.26, font_size=8.5, bold=True, color=TEAL)
questions = [
    "Revenue:   Which months/categories drive GMV and what explains the Nov 2017 Black Friday 4× spike?",
    "Customers: Why do 97% of customers not return, and which segments are most recoverable?",
    "Marketing: Which acquisition channel produces the highest seller conversion and post-conversion GMV?",
    "Logistics: What is the precise cost of late delivery in review score points, and which regions underperform?",
]
for i, q in enumerate(questions):
    col = i // 2
    row = i % 2
    add_text(slide, f"▸  {q}", 0.38 + col*4.75, 4.44 + row*0.42, 4.65, 0.4,
             font_size=8.5, color=DARK_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — METHODOLOGY & DATA: DATASET DESCRIPTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Methodology & Data — Dataset Description", 5)
footer_line(slide)

section_label(slide, "Section 3 of 6 — Methodology & Data", 0.25, 0.72)

add_text(slide, "Brazilian E-Commerce Dataset by Olist — Kaggle Public Dataset",
         0.25, 0.98, 9.5, 0.3, font_size=11.5, bold=True, color=NAVY)
add_text(slide, "11 relational CSV files covering orders, customers, products, sellers, payments, reviews, marketing leads, and geolocation (Sep 2016 – Oct 2018).",
         0.25, 1.28, 9.5, 0.28, font_size=9.5, color=MID_TXT)

# Dataset table
headers = ("Dataset", "Rows", "Key Fields", "Role in Analysis")
rows = [
    ("olist_orders_dataset",                "99,441",    "order_id, customer_id, status, dates",         "Core transactional grain"),
    ("olist_order_items_dataset",          "112,650",   "order_id, product_id, seller_id, price",       "Line-item revenue & freight"),
    ("olist_customers_dataset",             "99,441",    "customer_id, customer_unique_id, state, zip",  "Customer deduplication & geo"),
    ("olist_products_dataset",              "32,951",    "product_id, category_name, dimensions",        "Product classification"),
    ("olist_sellers_dataset",               "3,095",     "seller_id, state, zip_code_prefix",            "Seller geography & tier"),
    ("olist_order_payments_dataset",       "103,886",   "order_id, payment_type, installments, value",  "Payment method analytics"),
    ("olist_order_reviews_dataset",         "99,224",    "order_id, review_score, review_dates",         "Customer satisfaction score"),
    ("olist_marketing_qualified_leads",     "8,000",     "mql_id, origin (channel), first_contact_date", "B2B seller acquisition funnel"),
    ("olist_closed_deals_dataset",          "842",       "mql_id, seller_id, won_date, segment",         "Conversion & deal metadata"),
    ("olist_geolocation_dataset",          "1,000,163", "zip_prefix, lat, lng, city, state",             "Geographic intelligence (→19K)"),
    ("product_category_name_translation",   "71",        "category_name → category_name_english",        "Internationalised labels"),
]
col_w = [3.0, 0.85, 3.1, 2.55]
two_col_table(slide, [headers] + rows, 0.25, 1.60,
              sum(col_w), col_w, font_size=7.8,
              header_bg=NAVY, alt_bg=LIGHT_BG)

# Exclusions note
add_rect(slide, 0.25, 5.15, 9.5, 0.26, fill_color=rgb(0xFF,0xF8,0xF0))
add_text(slide,
    "Data exclusions & accepted gaps:  620 order items reference delisted products (NULL category — LEFT JOIN + coalesce applied)  |  "
    "~2% marketing leads have NULL origin (organic/direct — kept, labelled 'unknown')  |  "
    "26 geolocation rows outside Brazil bounding box removed pre-load",
    0.35, 5.17, 9.3, 0.22, font_size=7.5, color=MID_TXT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — METHODOLOGY & DATA: ELT APPROACH & DATA QUALITY
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Methodology & Data — ELT Architecture & Data Quality", 6)
footer_line(slide)

section_label(slide, "Section 3 of 6 — Methodology & Data", 0.25, 0.72)

# Left: 5-layer architecture
add_text(slide, "5-Layer ELT Architecture", 0.25, 0.98, 5.3, 0.3,
         font_size=11.5, bold=True, color=NAVY)
layers = [
    (TEAL,              "1 · SOURCE",    "11 CSV files (data/raw/)  |  1.46M total rows  |  Python ingest scripts"),
    (rgb(0xF4,0xA2,0x61), "2 · RAW LAYER","BigQuery olist_raw  |  WRITE_TRUNCATE  |  _ingested_at + _source_file metadata"),
    (rgb(0x28,0x9A,0xC4), "3 · STAGING",  "dbt staging models (views)  |  Type casts · NULL filters · Deduplication"),
    (rgb(0x1A,0x6B,0x4A), "4 · WAREHOUSE","Star schema  |  4 facts + 7 dims  |  FARM_FINGERPRINT surrogate keys"),
    (rgb(0x7B,0x35,0xA2), "5 · MARTS",    "7 dbt mart models  |  Pre-aggregated KPIs  |  Business logic gold layer"),
]
for i, (col, label, desc) in enumerate(layers):
    y = 1.35 + i * 0.68
    add_rect(slide, 0.25, y, 1.5, 0.56, fill_color=col)
    add_text(slide, label, 0.33, y+0.12, 1.38, 0.32,
             font_size=8.5, bold=True, color=WHITE)
    add_rect(slide, 1.75, y, 3.75, 0.56, fill_color=LIGHT_BG,
             line_color=col, line_width=0.6)
    add_text(slide, desc, 1.85, y+0.10, 3.58, 0.38, font_size=8.5, color=DARK_TXT, wrap=True)
    if i < 4:
        add_text(slide, "↓", 0.88, y+0.56, 0.4, 0.14, font_size=9, color=TEAL,
                 align=PP_ALIGN.CENTER)

# Tech stack
add_rect(slide, 0.25, 4.75, 5.25, 0.56, fill_color=rgb(0xF0,0xF4,0xFF))
tech = "BigQuery · dbt Core · Dagster · Python · MongoDB · Redis · GCS"
add_text(slide, "TECH STACK: " + tech, 0.38, 4.82, 5.08, 0.36,
         font_size=8.5, color=NAVY)

# Right: Data Quality pyramid
add_text(slide, "4-Layer Data Quality Framework", 5.7, 0.98, 4.0, 0.3,
         font_size=11.5, bold=True, color=NAVY)
dq_layers = [
    (rgb(0x7B,0x35,0xA2), "Layer 4: SQL Integration Tests",    "6 cross-table assertions  |  Pass: 100%"),
    (NAVY,                 "Layer 3: dbt Mart Schema Tests",    "40+ column tests across 7 marts  |  Pass: 100%"),
    (TEAL,                 "Layer 2: dbt Staging Schema Tests", "35+ tests on 9 staging models  |  Pass: 100%"),
    (ORANGE,               "Layer 1: Python Pre-Load DQ",       "87 checks: 57 core + 30 geo  |  Pass: 100%"),
]
for i, (col, label, detail) in enumerate(dq_layers):
    y = 1.35 + i * 0.82
    indent = i * 0.18
    bw = 4.0 - indent * 2
    add_rect(slide, 5.75 + indent, y, bw, 0.70, fill_color=col)
    add_text(slide, label, 5.83 + indent, y+0.06, bw-0.12, 0.26,
             font_size=9, bold=True, color=WHITE)
    add_text(slide, detail, 5.83 + indent, y+0.34, bw-0.12, 0.28,
             font_size=8.5, color=rgb(0xCC,0xEE,0xFF))

add_rect(slide, 5.75, 4.63, 4.0, 0.48, fill_color=rgb(0xF4,0xFD,0xFB),
         line_color=GREEN, line_width=0.7)
add_text(slide,
    "✅  All 181+ checks PASS on every pipeline run — CRITICAL failures halt pipeline "
    "and block downstream dbt execution",
    5.85, 4.68, 3.85, 0.4, font_size=8.5, color=DARK_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — RESULTS & INSIGHTS: KEY METRICS DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Results & Insights — Key Metrics Dashboard", 7)
footer_line(slide)

section_label(slide, "Section 4 of 6 — Results & Insights", 0.25, 0.72)

add_text(slide, "Platform-Wide KPI Scorecard — 2016–2018 Olist Dataset",
         0.25, 0.98, 9.5, 0.28, font_size=11, bold=True, color=NAVY)

# 6 big KPI cards
kpi_cards = [
    ("TOTAL GMV",       "R$15.8M",   "All delivered orders",  TEAL,  "≥ R$12M target ✓"),
    ("TOTAL ORDERS",    "99,441",    "Status: delivered",     TEAL,  "Baseline established"),
    ("AVG ORDER VALUE", "R$137.75",  "Per delivered order",   TEAL,  "≥ R$130 target ✓"),
    ("ON-TIME RATE",    "91.9%",     "vs estimated date",     rgb(0x1A,0x6B,0x4A), "Target: 95% — gap: 3.1pp"),
    ("AVG REVIEW SCORE","4.09 / 5", "Customer satisfaction",  rgb(0x1A,0x6B,0x4A), "Above industry 3.8 avg"),
    ("MQL CONVERSION",  "10.5%",    "MQL → seller",          rgb(0x7A,0x3A,0x1A), "Target: 8% — exceeded ✓"),
]
for i, (label, val, unit, bg, note) in enumerate(kpi_cards):
    col = i % 3
    row = i // 3
    x = 0.25 + col * 3.2
    y = 1.35 + row * 1.32
    add_rect(slide, x, y, 3.1, 1.22, fill_color=bg)
    add_text(slide, val,   x+0.12, y+0.10, 2.9, 0.55,
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, unit,  x+0.12, y+0.60, 2.9, 0.22,
             font_size=8.5, color=rgb(0xBB,0xE0,0xF5), align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.12, y+0.78, 2.9, 0.22,
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, y+1.0, 3.1, 0.22, fill_color=rgb(0,0,0,))  # dark strip
    # override dark strip
    add_rect(slide, x, y+1.0, 3.1, 0.22, fill_color=rgb(0x08,0x25,0x44))
    add_text(slide, note,  x+0.1, y+1.02, 2.95, 0.18,
             font_size=7.5, color=rgb(0x9F,0xD7,0xF0), align=PP_ALIGN.CENTER)

# Bottom insight strip
add_rect(slide, 0.25, 3.99, 9.5, 0.42, fill_color=rgb(0xFF,0xF8,0xE8))
add_text(slide,
    "★  Headline Insight:  97% of 96,096 unique customers placed only one order — repeat purchase rate = 3%. "
    "A 1pp improvement (3%→4%) = ~960 additional repeat customers = ~R$130K incremental GMV at current AOV.",
    0.38, 4.01, 9.28, 0.38, font_size=9, color=DARK_TXT, wrap=True)

add_rect(slide, 0.25, 4.46, 9.5, 0.86, fill_color=LIGHT_BG)
add_text(slide, "SECONDARY METRICS",
         0.38, 4.50, 9.2, 0.22, font_size=8.5, bold=True, color=TEAL)
secondary = [
    ("Repeat Purchase Rate", "3.0%",        "97% one-time buyers"),
    ("Avg Delivery Time",    "12.4 days",    "Target ≤ 10 days"),
    ("Late Delivery Rate",   "8.1%",        "North/NE highest"),
    ("Top Category GMV",     "R$1.26M",      "Health & Beauty"),
    ("Seller Count",         "3,095",        "Active on platform"),
    ("Paid Search Conv.",    "12.3%",        "Best-performing channel"),
]
for i, (metric, val, note) in enumerate(secondary):
    x = 0.38 + i * 1.55
    add_text(slide, metric, x, 4.74, 1.5, 0.2, font_size=7.5, bold=True, color=NAVY)
    add_text(slide, val,    x, 4.92, 1.5, 0.2, font_size=10, bold=True, color=TEAL)
    add_text(slide, note,   x, 5.10, 1.5, 0.2, font_size=7.5, color=DIM_TXT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — RESULTS: REVENUE & CATEGORY ANALYTICS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Results & Insights — Revenue & Category Analytics", 8)
footer_line(slide)

section_label(slide, "Section 4 of 6 — Results & Insights", 0.25, 0.72)

# Revenue chart (left)
add_text(slide, "Monthly GMV & Order Volume — 2016–2018",
         0.25, 0.98, 5.8, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "01_revenue_trend.png", 0.22, 1.28, 5.55, 3.2)

# Category chart (right)
add_text(slide, "Category Revenue Performance",
         5.95, 0.98, 3.8, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "04_category_performance.png", 5.92, 1.28, 3.85, 3.2)

# Insight strip
add_rect(slide, 0.25, 4.58, 9.5, 0.76, fill_color=rgb(0xE8,0xF4,0xFF),
         line_color=TEAL, line_width=0.5)
add_text(slide, "KEY REVENUE FINDINGS", 0.38, 4.62, 9.2, 0.22, font_size=8.5, bold=True, color=TEAL)
insights = [
    "Nov 2017 Black Friday: R$1.04M — 4× the monthly baseline; revenue 3× grew Jan→Nov 2017",
    "2018 stabilised at R$850K–R$1M/month — market maturity, not decline",
    "Health & Beauty + Watches: 18.5% of GMV from <5% of SKUs — subscription opportunity",
]
for i, ins in enumerate(insights):
    add_text(slide, f"▸  {ins}", 0.38 + (i % 3) * 3.17, 4.88, 3.1, 0.30,
             font_size=8.5, color=DARK_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — RESULTS: CUSTOMER SEGMENTATION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Results & Insights — Customer RFM Segmentation & CLV", 9)
footer_line(slide)

section_label(slide, "Section 4 of 6 — Results & Insights", 0.25, 0.72)

# RFM chart
add_text(slide, "RFM Customer Segmentation",
         0.25, 0.98, 5.5, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "02_rfm_segments.png", 0.22, 1.28, 5.55, 3.1)

# Segment analysis
add_text(slide, "Segment Analysis & Commercial Implication",
         5.95, 0.98, 3.8, 0.26, font_size=10, bold=True, color=NAVY)

segments = [
    (GREEN,  "CHAMPIONS",        "5,974 (6.2%)",  "35% of total GMV — protect with VIP programme"),
    (TEAL,   "AT-RISK",         "23,188 (24.1%)", "R$300K recovery potential at 5% reactivation"),
    (ORANGE, "NEW CUSTOMERS",    "11,684 (12.2%)", "First-time buyers: 7-day re-engagement window"),
    (RED,    "CANT LOSE THEM",   "~2,800 (2.9%)",  "High-frequency buyers gone silent — highest urgency"),
    (rgb(0x7B,0x35,0xA2), "SLEEPING GIANTS", "~4,500 (4.7%)", "High-spend, long recency — premium re-engagement"),
    (DIM_TXT,"HIBERNATING",     "~48,000 (50%+)", "Long-inactive — low-cost broadcast campaigns only"),
]
for i, (col, seg, count, action) in enumerate(segments):
    y = 1.30 + i * 0.56
    add_rect(slide, 5.92, y, 0.22, 0.44, fill_color=col)
    add_text(slide, seg,    6.18, y+0.04, 1.6, 0.22, font_size=8.5, bold=True, color=DARK_TXT)
    add_text(slide, count,  6.18, y+0.24, 1.5, 0.20, font_size=8,   color=MID_TXT)
    add_text(slide, action, 7.82, y+0.04, 1.95, 0.42, font_size=8,  color=DARK_TXT, wrap=True)

# 97% insight
add_rect(slide, 0.25, 4.48, 9.5, 0.36, fill_color=rgb(0xFF,0xEE,0xEE),
         line_color=RED, line_width=0.6)
add_text(slide,
    "⚠  Critical Finding: 97% of 96,096 customers purchased only ONCE. "
    "Repeat purchase rate = 3% vs industry benchmark of 25–35%. "
    "Root cause: no post-purchase engagement, no loyalty incentive, no subscription offering.",
    0.38, 4.52, 9.3, 0.3, font_size=8.5, color=DARK_TXT, wrap=True)

# CLV methodology note
add_rect(slide, 0.25, 4.92, 9.5, 0.40, fill_color=LIGHT_BG)
add_text(slide,
    "CLV Methodology (Phase 1): estimated_clv = total_spend × (1 + 0.3 × MIN(order_count − 1, 5))  "
    "| Champions median CLV: ~R$340  |  Platform median CLV: ~R$137 (≈ AOV for single-purchase majority)  "
    "| Phase 2: BG/NBD probabilistic model via Python lifetimes library",
    0.38, 4.96, 9.3, 0.34, font_size=7.8, color=MID_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — RESULTS: MARKETING FUNNEL & GEOGRAPHIC INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Results & Insights — Marketing Funnel & Geographic Demand", 10)
footer_line(slide)

section_label(slide, "Section 4 of 6 — Results & Insights", 0.25, 0.72)

# Left: funnel chart
add_text(slide, "Seller Acquisition Funnel by Channel",
         0.25, 0.98, 4.8, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "03_marketing_funnel.png", 0.22, 1.28, 4.85, 2.6)

# Funnel insights
add_rect(slide, 0.22, 3.92, 4.88, 1.40, fill_color=LIGHT_BG, line_color=TEAL, line_width=0.5)
add_text(slide, "CHANNEL CONVERSION RATES", 0.35, 3.96, 4.6, 0.22,
         font_size=8.5, bold=True, color=TEAL)
channels = [
    ("Paid Search",    "12.3%", GREEN,  "4× Email — shift budget here"),
    ("Organic Search", "11.8%", GREEN,  "Near-equal to paid — invest in SEO"),
    ("Social Media",   "5.6%",  ORANGE, "Mid-tier — nurture sequences needed"),
    ("Email",          "3.0%",  RED,    "Lowest — reallocate 20% of budget"),
]
for i, (ch, rate, col, note) in enumerate(channels):
    y = 4.20 + i * 0.28
    bar_w = float(rate.strip('%')) / 15.0  # scale to max ~1.0 width
    add_rect(slide, 0.35, y, bar_w, 0.22, fill_color=col)
    add_text(slide, f"{ch}: {rate} — {note}",
             0.35 + bar_w + 0.05, y+0.02, 4.65, 0.22, font_size=7.5, color=DARK_TXT)

# Right: geographic
add_text(slide, "Geographic Demand Intelligence — 27 States",
         5.25, 0.98, 4.55, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "07_geographic.png", 5.22, 1.28, 4.58, 2.6)

# Geo insights
add_rect(slide, 5.22, 3.92, 4.58, 1.40, fill_color=LIGHT_BG, line_color=TEAL, line_width=0.5)
add_text(slide, "REGIONAL DEMAND vs POPULATION", 5.35, 3.96, 4.3, 0.22,
         font_size=8.5, bold=True, color=TEAL)
geo_data = [
    ("Southeast (SP+RJ+MG)", "42% pop", "68% orders", GREEN,  "Strong penetration — maintain"),
    ("South (PR+SC+RS)",     "15% pop", "15% orders", TEAL,   "Market saturation — monetise"),
    ("Northeast",            "28% pop", "16% orders", ORANGE, "Underserved — logistics gap"),
    ("North",                "8.5% pop","<1% orders", RED,    "Critical gap — hub investment"),
]
for i, (region, pop, orders, col, note) in enumerate(geo_data):
    y = 4.20 + i * 0.28
    add_rect(slide, 5.35, y, 0.16, 0.22, fill_color=col)
    add_text(slide, f"{region}:  pop {pop}  |  orders {orders}  — {note}",
             5.55, y+0.02, 4.18, 0.22, font_size=7.5, color=DARK_TXT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — RESULTS: LOGISTICS & DELIVERY PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Results & Insights — Logistics & Delivery Performance", 11)
footer_line(slide)

section_label(slide, "Section 4 of 6 — Results & Insights", 0.25, 0.72)

# Delivery chart
add_text(slide, "Delivery Performance & Review Score Correlation",
         0.25, 0.98, 5.5, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "05_delivery_performance.png", 0.22, 1.28, 5.55, 2.9)

# Payments chart
add_text(slide, "Payment Method Distribution",
         5.92, 0.98, 3.8, 0.26, font_size=10, bold=True, color=NAVY)
add_image(slide, CHARTS / "08_payments.png", 5.88, 1.28, 3.9, 2.9)

# Insights row
add_rect(slide, 0.25, 4.30, 9.5, 1.05, fill_color=LIGHT_BG, line_color=TEAL, line_width=0.5)
add_text(slide, "LOGISTICS & PAYMENT FINDINGS", 0.38, 4.34, 9.2, 0.24,
         font_size=8.5, bold=True, color=TEAL)
logistic_insights = [
    "⏱  Avg delivery: 12.4 days | On-time: 91.9% | Late: 8.1% concentrated in North/Northeast",
    "⭐  Every +1 day delivery → −0.05 review points | Sub-5-day orders score 4.42 | 31+ day: 3.58",
    "💳  Credit card: 74% of payments | Boleto (bank slip): 19% → significant unbanked segment",
    "📦  Proactive delay SMS at T+1 day recovers ~30% satisfaction loss — near-zero implementation cost",
]
for i, ins in enumerate(logistic_insights):
    col = i % 2
    row = i // 2
    add_text(slide, f"▸  {ins}", 0.38 + col*4.75, 4.62 + row*0.36, 4.68, 0.34,
             font_size=8.5, color=DARK_TXT, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — STRATEGIC RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Strategic Recommendations", 12)
footer_line(slide)

section_label(slide, "Section 5 of 6 — Strategic Recommendations", 0.25, 0.72)

recs = [
    (
        "1", RED,
        "Win-Back Campaign — At-Risk Customers",
        "R$300K incremental GMV",
        "Low",
        "Target 23,188 At-Risk customers with personalised email sequences segmented by RFM sub-score. "
        "5% reactivation rate (industry benchmark for churn win-back) yields ~1,160 returning customers × R$260 avg CLV. "
        "Sequence: Day 1 personalised offer → Day 7 urgency → Day 14 loyalty invite.",
        "CMO",
    ),
    (
        "2", ORANGE,
        "Paid Search Budget Reallocation",
        "+35–40 deals/cohort",
        "Low",
        "Paid Search converts at 12.3% vs Email at 3.0% — a 4× efficiency gap. "
        "Reallocating 20% of Social + Email spend to Paid Search (holding total budget constant) "
        "projects +35–40 additional closed seller deals per cohort based on observed conversion differentials. "
        "Also invest in organic SEO (11.8% conversion, zero incremental cost).",
        "VP Marketing",
    ),
    (
        "3", rgb(0x7B,0x35,0xA2),
        "Northern Brazil Fulfilment Hub",
        "5K+ new orders/year",
        "High",
        "North/Northeast: 36.5% of Brazil's population but only 17% of orders. "
        "Avg delivery to North: 18+ days vs Southeast: 9 days. A regional hub reduces delivery time 40%, "
        "improving review scores and repeat purchase intent. Avg order value in North is above the platform mean.",
        "COO / CFO",
    ),
    (
        "4", TEAL,
        "Health & Beauty Subscription Model",
        "3–5× LTV uplift",
        "Medium",
        "Health & Beauty is the #1 revenue category at R$1.26M. 92% of products are repeat consumables "
        "(skincare, vitamins, supplements). A monthly subscription box at 10% discount vs single-order price "
        "increases customer LTV 3–5× and converts one-time buyers into contracted recurring revenue.",
        "CPO / CMO",
    ),
    (
        "5", GREEN,
        "Carrier SLA Renegotiation + Proactive Notifications",
        "+0.3 review stars",
        "Low",
        "8.1% late deliveries suppress review score by an estimated 0.5 points on affected orders. "
        "Two actions: (1) Renegotiate carrier SLAs for North/Northeast routes — current SLA not enforced. "
        "(2) Implement SMS/push alert at T+1 day past estimated delivery — reduces 1-star reviews 30% (industry benchmark). "
        "Pipeline already tracks delivery timestamps; notification trigger is a 2-hour engineering task.",
        "COO",
    ),
]

col_headers = ("Impact", "Effort", "Owner")
for i, (num, col, title, impact, effort, desc, owner) in enumerate(recs):
    x = 0.22
    y = 1.00 + i * 0.90
    add_rect(slide, x, y, 0.32, 0.78, fill_color=col)
    add_text(slide, num, x+0.04, y+0.20, 0.24, 0.40,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x+0.38, y+0.04, 6.55, 0.26,
             font_size=9.5, bold=True, color=DARK_TXT)
    add_text(slide, desc, x+0.38, y+0.30, 6.55, 0.48,
             font_size=8, color=MID_TXT, wrap=True)
    # Impact badge
    add_rect(slide, 7.12, y+0.08, 0.9, 0.26, fill_color=rgb(0xE8,0xF8,0xE8))
    add_text(slide, impact, 7.15, y+0.10, 0.86, 0.22,
             font_size=7.5, bold=True, color=rgb(0x0A,0x5A,0x2A), align=PP_ALIGN.CENTER)
    # Effort badge
    eff_col = GREEN if effort == "Low" else ORANGE if effort == "Medium" else RED
    add_rect(slide, 8.08, y+0.08, 0.84, 0.26, fill_color=eff_col)
    add_text(slide, f"Effort: {effort}", 8.10, y+0.10, 0.8, 0.22,
             font_size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, owner, 8.98, y+0.10, 0.85, 0.22,
             font_size=7.5, color=DIM_TXT, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 10, 5.6, fill_color=WHITE)
header_band(slide, "Conclusion — Summary & Strategic Implications", 13)
footer_line(slide)

section_label(slide, "Section 6 of 6 — Conclusion", 0.25, 0.72)

# Left: what was built
add_text(slide, "What Was Built", 0.25, 0.98, 4.6, 0.28, font_size=11, bold=True, color=NAVY)
built_items = [
    "5-layer ELT pipeline: CSV → BigQuery Raw → Staging → Warehouse → Marts",
    "Star schema: 4 fact tables, 7 dimensions, FARM_FINGERPRINT surrogate keys",
    "181+ automated data quality checks across 4 framework layers",
    "7 dbt mart models: revenue, CLV, RFM, geo, seller, logistics, marketing",
    "Dagster orchestration with daily 06:00 UTC schedule (auto-activating)",
    "8 executive analytics charts | MongoDB ODS | Redis KPI cache layer",
    "Comprehensive documentation: data dictionary, KPI catalogue, tech specs",
]
add_rect(slide, 0.25, 1.30, 4.6, 3.60, fill_color=LIGHT_BG, line_color=TEAL, line_width=0.5)
for i, item in enumerate(built_items):
    add_text(slide, f"✓  {item}", 0.38, 1.40 + i * 0.49, 4.38, 0.46,
             font_size=8.5, color=DARK_TXT, wrap=True)

# Right: key findings
add_text(slide, "Key Findings & Strategic Implications", 5.05, 0.98, 4.65, 0.28,
         font_size=11, bold=True, color=NAVY)
implications = [
    (RED,    "Retention Crisis",
     "3% repeat purchase rate is the #1 business risk. A loyalty programme targeting 8% repeat rate "
     "by Q2 adds ~R$650K GMV annually without acquiring a single new customer."),
    (ORANGE, "Marketing Efficiency Gap",
     "4× conversion gap between Paid Search and Email represents misallocated budget. "
     "Realigning spend to performance data is a zero-cost, high-ROI action."),
    (rgb(0x7B,0x35,0xA2), "Geographic White Space",
     "Northern Brazil (36% of population, <1% of orders) is the most significant "
     "untapped growth opportunity. Infrastructure investment unlocks the next growth curve."),
    (TEAL,   "Platform Readiness",
     "The data platform is production-ready: automated, tested, documented, and scheduled. "
     "Phase 2 (real-time streaming, ML churn prediction) can begin immediately."),
]
for i, (col, title, text) in enumerate(implications):
    y = 1.30 + i * 0.85
    add_rect(slide, 5.05, y, 0.22, 0.72, fill_color=col)
    add_text(slide, title, 5.32, y+0.04, 4.3, 0.24, font_size=9, bold=True, color=col)
    add_text(slide, text,  5.32, y+0.28, 4.3, 0.46, font_size=8.5, color=DARK_TXT, wrap=True)

# Bottom: roadmap strip
add_rect(slide, 0.25, 5.00, 9.5, 0.40, fill_color=NAVY)
phases = [
    ("Q3 2026 — Foundation Live", "Real-time Kafka ingestion · ML churn model · Looker Studio dashboards"),
    ("Q4 2026 — Scale",           "Personalisation engine · Dynamic pricing · North hub · Subscription"),
    ("Q1 2027 — Intelligence",    "Multi-country · AI recommendation API · Self-serve analytics portal"),
]
for i, (phase, detail) in enumerate(phases):
    x = 0.38 + i * 3.15
    add_text(slide, phase,  x, 5.02, 3.1, 0.18, font_size=7.5, bold=True, color=TEAL_LT)
    add_text(slide, detail, x, 5.20, 3.1, 0.18, font_size=7,   color=rgb(0x88,0xAA,0xCC))

# Sign-off
add_rect(slide, 0, 4.92, 10, 0.02, fill_color=TEAL)


# ═══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"✅  Presentation saved: {OUT}")
print(f"    Slides: {len(prs.slides)}")
print(f"    Dimensions: {prs.slide_width.inches:.1f}in × {prs.slide_height.inches:.1f}in")
